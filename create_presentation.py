import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("Error: python-pptx library not found.")
    print("Please run: pip install python-pptx")
    sys.exit(1)

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 Widescreen standard
prs.slide_height = Inches(7.5)

# Color Palette (Premium Indigo / Slate Theme)
DARK_BG = RGBColor(15, 23, 42)      # Slate 900
TEXT_LIGHT = RGBColor(248, 250, 252) # Slate 50
TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
ACCENT_VIOLET = RGBColor(139, 92, 246) # Violet 500
ACCENT_BLUE = RGBColor(56, 189, 248)    # Sky 400

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_header(slide, title_text, category_text="VYBE DEVOPS"):
    # Category / Tag
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_text.upper()
    p_tag.font.name = 'Segoe UI'
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_VIOLET
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = 'Segoe UI'
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT

blank_layout = prs.slide_layouts[6]

# ==========================================
# Slide 1: Project Overview & Objectives
# ==========================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

# Large Main Title Box
title_box = slide1.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.2))
tf = title_box.text_frame
tf.word_wrap = True

p_main = tf.paragraphs[0]
p_main.text = "Vybe Social Media Platform"
p_main.font.name = 'Segoe UI'
p_main.font.size = Pt(44)
p_main.font.bold = True
p_main.font.color.rgb = TEXT_LIGHT

p_sub = tf.add_paragraph()
p_sub.text = "Enterprise-Grade Cloud-Native Deployment"
p_sub.font.name = 'Segoe UI'
p_sub.font.size = Pt(22)
p_sub.font.bold = True
p_sub.font.color.rgb = ACCENT_VIOLET
p_sub.space_before = Pt(8)

# Core Goal / Description Box
goal_box = slide1.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.5))
tf_goal = goal_box.text_frame
tf_goal.word_wrap = True

p_goal_lbl = tf_goal.paragraphs[0]
p_goal_lbl.text = "PROJECT OBJECTIVE & GOAL:"
p_goal_lbl.font.name = 'Segoe UI'
p_goal_lbl.font.size = Pt(14)
p_goal_lbl.font.bold = True
p_goal_lbl.font.color.rgb = ACCENT_BLUE

p_goal_desc = tf_goal.add_paragraph()
p_goal_desc.text = "Migrate from a local containerized development application to a high-availability, fully automated, and monitored AWS Kubernetes environment. The objective is to build a reliable infrastructure capable of scaling real-time chat socket connections, media uploads, and observability feeds globally."
p_goal_desc.font.name = 'Segoe UI'
p_goal_desc.font.size = Pt(16)
p_goal_desc.font.color.rgb = TEXT_MUTED
p_goal_desc.space_before = Pt(8)

p_author = tf_goal.add_paragraph()
p_author.text = "Presenter: Satish Raut"
p_author.font.name = 'Segoe UI'
p_author.font.size = Pt(14)
p_author.font.bold = True
p_author.font.color.rgb = TEXT_LIGHT
p_author.space_before = Pt(25)


# ==========================================
# Slide 2: Infrastructure as Code (Terraform)
# ==========================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "Infrastructure as Code (IaC) via Terraform", "Slide 2")

content_box2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf2 = content_box2.text_frame
tf2.word_wrap = True

tf2.paragraphs[0].text = "Core Infrastructure Components:"
tf2.paragraphs[0].font.name = 'Segoe UI'
tf2.paragraphs[0].font.size = Pt(18)
tf2.paragraphs[0].font.bold = True
tf2.paragraphs[0].font.color.rgb = ACCENT_BLUE

tf2.paragraphs[0].space_after = Pt(10)

points2 = [
    ("Custom AWS VPC", "Provisioned isolated public and private subnets across multiple availability zones (us-east-1a, us-east-1b) for secure high-availability load balancing."),
    ("AWS EKS Cluster", "Created a managed Kubernetes deployment cluster using t3.small EC2 worker nodes to host container workloads with IAM security roles."),
    ("Amazon ECR Registries", "Set up private Elastic Container Registries (ECR) for the React frontend and Node.js backend to host built Docker images.")
]

for title, desc in points2:
    p_t = tf2.add_paragraph()
    p_t.text = f"•  {title}"
    p_t.font.name = 'Segoe UI'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_LIGHT
    p_t.space_before = Pt(12)
    
    p_d = tf2.add_paragraph()
    p_d.text = f"   {desc}"
    p_d.font.name = 'Segoe UI'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = TEXT_MUTED
    p_d.space_before = Pt(3)

p_val2 = tf2.add_paragraph()
p_val2.text = "Value: Standardized, reproducible infrastructure built cleanly in 15 minutes."
p_val2.font.name = 'Segoe UI'
p_val2.font.size = Pt(15)
p_val2.font.bold = True
p_val2.font.color.rgb = ACCENT_VIOLET
p_val2.space_before = Pt(25)


# ==========================================
# Slide 3: Automated CI/CD Pipeline (Jenkins)
# ==========================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "Automated Continuous Integration & Deployment", "Slide 3")

content_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf3 = content_box3.text_frame
tf3.word_wrap = True

tf3.paragraphs[0].text = "CI/CD Pipeline Flow:"
tf3.paragraphs[0].font.name = 'Segoe UI'
tf3.paragraphs[0].font.size = Pt(18)
tf3.paragraphs[0].font.bold = True
tf3.paragraphs[0].font.color.rgb = ACCENT_BLUE
tf3.paragraphs[0].space_after = Pt(10)

points3 = [
    ("GitHub Webhook Trigger", "Pushing code to the main branch automatically calls our dedicated Jenkins EC2 server to trigger a fresh build sequence."),
    ("Multi-Stage Docker Builds", "Jenkins runs isolated, lightweight Docker builds, compiling the React client inside Nginx and setting up Node.js server runtimes."),
    ("Rolling Updates to AWS EKS", "Tags fresh images, pushes them to AWS ECR, and applies Kubernetes manifests to EKS with zero deployment downtime.")
]

for title, desc in points3:
    p_t = tf3.add_paragraph()
    p_t.text = f"•  {title}"
    p_t.font.name = 'Segoe UI'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_LIGHT
    p_t.space_before = Pt(12)
    
    p_d = tf3.add_paragraph()
    p_d.text = f"   {desc}"
    p_d.font.name = 'Segoe UI'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = TEXT_MUTED
    p_d.space_before = Pt(3)

p_val3 = tf3.add_paragraph()
p_val3.text = "Value: True DevOps continuous integration with zero manual steps."
p_val3.font.name = 'Segoe UI'
p_val3.font.size = Pt(15)
p_val3.font.bold = True
p_val3.font.color.rgb = ACCENT_VIOLET
p_val3.space_before = Pt(25)


# ==========================================
# Slide 4: Advanced Performance Tuning (The "Pro" Slide)
# ==========================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_background(slide4)
add_header(slide4, "Resolving Real-World Network & Timeout Limits", "Slide 4")

content_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf4 = content_box4.text_frame
tf4.word_wrap = True

tf4.paragraphs[0].text = "Tuning Network & Connection Proxies:"
tf4.paragraphs[0].font.name = 'Segoe UI'
tf4.paragraphs[0].font.size = Pt(18)
tf4.paragraphs[0].font.bold = True
tf4.paragraphs[0].font.color.rgb = ACCENT_BLUE
tf4.paragraphs[0].space_after = Pt(10)

points4 = [
    ("Disabling Ingress Proxy Buffering", "Disabled Nginx response buffering via ingress annotations to allow Server-Sent Events (SSE) chat streams to flow instantly to clients without delay."),
    ("SSE Keep-Alive Heartbeat (Timeout Prevention)", "Implemented a 20-second backend keep-alive heartbeat ping to prevent AWS Load Balancers from dropping idle TCP connections, fixing ERR_INCOMPLETE_CHUNKED_ENCODING errors."),
    ("Nginx Client Body Limit Expansion", "Configured 'proxy-body-size: 50m' annotation in the Ingress manifest to allow large media uploads, bypassing the default 1MB restriction.")
]

for title, desc in points4:
    p_t = tf4.add_paragraph()
    p_t.text = f"•  {title}"
    p_t.font.name = 'Segoe UI'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_LIGHT
    p_t.space_before = Pt(12)
    
    p_d = tf4.add_paragraph()
    p_d.text = f"   {desc}"
    p_d.font.name = 'Segoe UI'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = TEXT_MUTED
    p_d.space_before = Pt(3)

p_val4 = tf4.add_paragraph()
p_val4.text = "Value: Custom proxy tuning to match real-world distributed messaging states."
p_val4.font.name = 'Segoe UI'
p_val4.font.size = Pt(15)
p_val4.font.bold = True
p_val4.font.color.rgb = ACCENT_VIOLET
p_val4.space_before = Pt(25)


# ==========================================
# Slide 5: Observability & Monitoring (Prometheus & Grafana)
# ==========================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "Cluster Observability & Metrics Infrastructure", "Slide 5")

content_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf5 = content_box5.text_frame
tf5.word_wrap = True

tf5.paragraphs[0].text = "Monitoring Stack Architecture:"
tf5.paragraphs[0].font.name = 'Segoe UI'
tf5.paragraphs[0].font.size = Pt(18)
tf5.paragraphs[0].font.bold = True
tf5.paragraphs[0].font.color.rgb = ACCENT_BLUE
tf5.paragraphs[0].space_after = Pt(10)

points5 = [
    ("Helm Deployment Stack", "Deployed the community-verified 'kube-prometheus-stack' Helm Chart inside a dedicated namespace, spinning up node exporters, state metrics, and the visual dashboards."),
    ("Data Collection and Scrapes", "Prometheus dynamically gathers CPU limits, memory quotas, and network inputs from every pod, storing them as time-series metrics."),
    ("Grafana Analytics Console", "Configured interactive visual dashboards linked to Prometheus data-sources, allowing operators to monitor memory limits and CPU utilization in real-time.")
]

for title, desc in points5:
    p_t = tf5.add_paragraph()
    p_t.text = f"•  {title}"
    p_t.font.name = 'Segoe UI'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_LIGHT
    p_t.space_before = Pt(12)
    
    p_d = tf5.add_paragraph()
    p_d.text = f"   {desc}"
    p_d.font.name = 'Segoe UI'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = TEXT_MUTED
    p_d.space_before = Pt(3)

p_val5 = tf5.add_paragraph()
p_val5.text = "Value: Enterprise-grade cluster health reporting."
p_val5.font.name = 'Segoe UI'
p_val5.font.size = Pt(15)
p_val5.font.bold = True
p_val5.font.color.rgb = ACCENT_VIOLET
p_val5.space_before = Pt(25)


# ==========================================
# Slide 6: Live Demonstration & Q&A
# ==========================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "Production Deliverables & Live Demo", "Slide 6")

content_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf6 = content_box6.text_frame
tf6.word_wrap = True

tf6.paragraphs[0].text = "Live Staging Environment Details:"
tf6.paragraphs[0].font.name = 'Segoe UI'
tf6.paragraphs[0].font.size = Pt(18)
tf6.paragraphs[0].font.bold = True
tf6.paragraphs[0].font.color.rgb = ACCENT_BLUE
tf6.paragraphs[0].space_after = Pt(10)

points6 = [
    ("Live Global Website", "Staging site hosted at: http://www.bestvybe.live (Connected via CNAME pointing to our AWS Elastic Load Balancer)."),
    ("Observability Link", "Grafana dashboards running at http://localhost:3100 (Connected via secure local port-forwarding)."),
    ("Clean Resource Teardown", "Execution of 'terraform destroy' will cleanly decommission the entire infrastructure when staging presentation is complete.")
]

for title, desc in points6:
    p_t = tf6.add_paragraph()
    p_t.text = f"•  {title}"
    p_t.font.name = 'Segoe UI'
    p_t.font.size = Pt(16)
    p_t.font.bold = True
    p_t.font.color.rgb = TEXT_LIGHT
    p_t.space_before = Pt(12)
    
    p_d = tf6.add_paragraph()
    p_d.text = f"   {desc}"
    p_d.font.name = 'Segoe UI'
    p_d.font.size = Pt(14)
    p_d.font.color.rgb = TEXT_MUTED
    p_d.space_before = Pt(3)

p_val6 = tf6.add_paragraph()
p_val6.text = "Thank you! Open for Q&A."
p_val6.font.name = 'Segoe UI'
p_val6.font.size = Pt(18)
p_val6.font.bold = True
p_val6.font.color.rgb = ACCENT_VIOLET
p_val6.space_before = Pt(30)

# Save the presentation file
output_path = "Vybe_DevOps_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved successfully to {os.path.abspath(output_path)}")
